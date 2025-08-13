import fitz  # PyMuPDF for PDF parsing
from pptx import Presentation
import os

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    doc = fitz.open(pdf_path)
    text = "\n".join(page.get_text("text") for page in doc)
    return text if text else "No text found in PDF."

def extract_text_from_ppt(ppt_path):
    """Extract text from a PowerPoint file."""
    prs = Presentation(ppt_path)
    text = "\n".join(
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in prs.slides
    )
    return text if text else "No text found in PPT."

def extract_text(file_path):
    """Determine file type and extract text accordingly."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_text_from_ppt(file_path)
    return "Unsupported file format."
